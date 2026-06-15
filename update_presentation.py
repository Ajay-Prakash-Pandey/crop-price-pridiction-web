from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUTPUT = Path("CropSmart_BCA_Final_Project_Viva.pptx")

NAVY = RGBColor(13, 32, 45)
GREEN = RGBColor(35, 133, 91)
LIGHT_GREEN = RGBColor(226, 244, 235)
GOLD = RGBColor(244, 180, 62)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(31, 46, 55)
MUTED = RGBColor(91, 108, 116)
PALE = RGBColor(246, 249, 247)


def add_text(slide, x, y, w, h, text, size=20, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", margin=0.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title(slide, title, number):
    add_rect(slide, 0, 0, 13.333, 0.17, GREEN, radius=False)
    add_text(slide, 0.55, 0.34, 11.7, 0.62, title, 28, NAVY, True)
    add_text(slide, 12.1, 0.38, 0.7, 0.4, f"{number:02}", 14, GREEN, True,
             PP_ALIGN.RIGHT)
    add_text(slide, 0.55, 7.1, 12.2, 0.22,
             "CropSmart  |  BCA Final-Year Project  |  June 16, 2026",
             9, MUTED)


def add_card(slide, x, y, w, h, heading, body, accent=GREEN):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(218, 229, 223))
    add_rect(slide, x, y, 0.08, h, accent, radius=False)
    add_text(slide, x + 0.2, y + 0.13, w - 0.35, 0.35, heading, 16, accent, True)
    add_text(slide, x + 0.2, y + 0.53, w - 0.35, h - 0.65, body, 12, TEXT)


def add_bullets(slide, x, y, w, h, items, size=18, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(10)
        paragraph.level = 0
        paragraph.text = f"•  {item}"
    return box


def base_slide(prs, title, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PALE
    add_title(slide, title, number)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_rect(slide, 0, 0, 0.18, 7.5, GREEN, radius=False)
    add_text(slide, 0.9, 0.8, 3.2, 0.4, "BCA FINAL-YEAR PROJECT", 14, GOLD, True)
    add_text(slide, 0.9, 1.35, 7.5, 1.0, "CropSmart", 48, WHITE, True)
    add_text(slide, 0.9, 2.25, 8.8, 0.9,
             "Crop Price Prediction and\nMarket Intelligence System",
             27, LIGHT_GREEN, True)
    add_rect(slide, 9.8, 1.15, 2.2, 2.2, GREEN)
    add_text(slide, 10.1, 1.52, 1.6, 0.6, "CS", 39, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 10.05, 2.25, 1.7, 0.55, "SMART\nFARMING", 12, WHITE, True,
             PP_ALIGN.CENTER)
    add_text(slide, 0.9, 4.55, 5.8, 0.45,
             "Presented by: __________________________", 16, WHITE)
    add_text(slide, 0.9, 5.05, 5.8, 0.45,
             "Roll No.: ______________________________", 16, WHITE)
    add_text(slide, 0.9, 5.55, 5.8, 0.45,
             "College: _______________________________", 16, WHITE)
    add_text(slide, 0.9, 6.45, 7.5, 0.35,
             "Final Project Viva  |  June 16, 2026", 14, GOLD, True)

    # 2. Problem
    slide = base_slide(prs, "Problem Statement", 2)
    add_text(slide, 0.7, 1.15, 5.6, 0.55,
             "Crop prices are difficult to track and predict.", 25, NAVY, True)
    add_bullets(slide, 0.72, 1.85, 5.7, 3.9, [
        "Prices change due to weather, supply, demand, inflation, transport cost, policy, and global risk.",
        "Farmers and buyers often check multiple sources manually.",
        "Local and global market information is not available in one simple dashboard.",
        "Future prices are uncertain, making planning and selling decisions difficult.",
    ], 17)
    add_rect(slide, 7.0, 1.35, 5.2, 4.7, LIGHT_GREEN)
    add_text(slide, 7.45, 1.75, 4.3, 0.55, "CropSmart Solution", 25, GREEN, True,
             PP_ALIGN.CENTER)
    add_text(slide, 7.55, 2.55, 4.1, 2.7,
             "One browser-based system for:\n\n"
             "CURRENT PRICE\nHISTORICAL TREND\n24-MONTH FORECAST\nPRICE DRIVERS\nCSV / JSON EXPORT",
             18, NAVY, True, PP_ALIGN.CENTER)

    # 3. Objectives
    slide = base_slide(prs, "Project Objectives", 3)
    cards = [
        ("Current Price", "Show a direct public benchmark quote when available, otherwise a labelled estimate."),
        ("Future Forecast", "Predict selected periods and produce a complete 24-month monthly outlook."),
        ("Historical View", "Display up to 10 years of observed benchmark history for supported crops."),
        ("Explainability", "Show the market factors and sources influencing the expected price."),
        ("Global Options", "Support 22 crops or variants and 14 major regions with currency conversion."),
        ("Export", "Download current, historical, forecast, source, and model details as CSV or JSON."),
    ]
    positions = [(0.65, 1.3), (4.55, 1.3), (8.45, 1.3),
                 (0.65, 3.85), (4.55, 3.85), (8.45, 3.85)]
    for (heading, body), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 3.55, 1.95, heading, body)

    # 4. Technologies
    slide = base_slide(prs, "Technology Stack", 4)
    stack = [
        ("Frontend", "HTML5\nCSS3\nJavaScript", GREEN),
        ("Backend", "Python\nHTTP server\nREST-style API", RGBColor(39, 105, 156)),
        ("Machine Learning", "Ridge regression\nTime-series lags\nHoldout validation", RGBColor(141, 91, 166)),
        ("Data", "Public web APIs\nJSON cache\nCSV datasets", RGBColor(203, 126, 36)),
    ]
    for index, (heading, body, color) in enumerate(stack):
        x = 0.65 + index * 3.15
        add_rect(slide, x, 1.45, 2.75, 4.7, WHITE, line=RGBColor(218, 229, 223))
        add_rect(slide, x, 1.45, 2.75, 0.8, color)
        add_text(slide, x + 0.1, 1.58, 2.55, 0.45, heading, 18, WHITE, True,
                 PP_ALIGN.CENTER)
        add_text(slide, x + 0.25, 2.7, 2.25, 2.5, body, 19, NAVY, True,
                 PP_ALIGN.CENTER)

    # 5. Architecture
    slide = base_slide(prs, "System Architecture", 5)
    nodes = [
        (0.55, "USER", "Selects crop,\nregion, currency,\nand period"),
        (3.05, "FRONTEND", "HTML, CSS,\nand JavaScript\ndashboard"),
        (5.55, "PYTHON API", "Prediction,\noptions, health,\nand export"),
        (8.05, "DATA + MODEL", "Public signals,\nhistory, cache,\nand ridge model"),
        (10.55, "OUTPUT", "Price, charts,\ndrivers, MAPE,\nand export"),
    ]
    for index, (x, heading, body) in enumerate(nodes):
        add_rect(slide, x, 2.0, 2.0, 2.6, WHITE, line=GREEN)
        add_text(slide, x + 0.12, 2.25, 1.76, 0.45, heading, 15, GREEN, True,
                 PP_ALIGN.CENTER)
        add_text(slide, x + 0.18, 2.95, 1.64, 1.25, body, 14, TEXT, False,
                 PP_ALIGN.CENTER)
        if index < len(nodes) - 1:
            add_text(slide, x + 2.02, 2.85, 0.45, 0.5, "→", 27, GOLD, True,
                     PP_ALIGN.CENTER)
    add_text(slide, 1.0, 5.25, 11.3, 0.65,
             "Request flow: Browser → API v2 → live data and history → forecast model → dashboard response",
             18, NAVY, True, PP_ALIGN.CENTER)

    # 6. Data
    slide = base_slide(prs, "Data Sources and Processing", 6)
    sources = [
        ("Market Data", "Yahoo Finance benchmark quotes and monthly futures closes."),
        ("Economic Data", "World Bank inflation and GDP indicators."),
        ("Weather Data", "Open-Meteo temperature and rainfall signals."),
        ("Risk Data", "GDELT conflict and unrest news signals."),
    ]
    for i, (heading, body) in enumerate(sources):
        add_card(slide, 0.65 + (i % 2) * 6.15, 1.25 + (i // 2) * 2.05,
                 5.7, 1.55, heading, body)
    add_rect(slide, 0.65, 5.45, 11.85, 0.8, NAVY)
    add_text(slide, 0.9, 5.58, 11.35, 0.48,
             "CACHE: recent API responses are stored in JSON to improve speed and handle temporary failures.",
             15, WHITE, True, PP_ALIGN.CENTER)

    # 7. Model
    slide = base_slide(prs, "Prediction Method", 7)
    add_card(slide, 0.65, 1.25, 5.75, 4.8, "Path A: Observed Benchmark Series",
             "1. Load at least 36 monthly prices.\n\n"
             "2. Build trend, seasonal sine/cosine, and 1-, 3-, and 12-month lag features.\n\n"
             "3. Train ridge autoregression on older observations.\n\n"
             "4. Test on the latest 6-12 months.\n\n"
             "5. Forecast the next 24 months recursively.")
    add_card(slide, 6.9, 1.25, 5.75, 4.8, "Path B: Scenario Fallback",
             "Used when a crop has no suitable public historical series.\n\n"
             "The ridge model uses crop, region, inflation, demand, fuel cost, climate, policy, risk, supply, and seasonality features.\n\n"
             "The interface clearly labels this output as a scenario-model fallback.",
             RGBColor(203, 126, 36))

    # 8. Validation
    slide = base_slide(prs, "Model Validation", 8)
    add_rect(slide, 0.7, 1.25, 3.65, 4.9, LIGHT_GREEN)
    add_text(slide, 1.05, 1.65, 2.95, 0.55, "Chronological Holdout", 22, GREEN, True,
             PP_ALIGN.CENTER)
    add_text(slide, 1.0, 2.45, 3.05, 2.7,
             "TRAIN\nOlder monthly prices\n\n↓\n\nTEST\nLatest 6-12 months",
             20, NAVY, True, PP_ALIGN.CENTER)
    add_card(slide, 4.8, 1.25, 3.55, 2.15, "MAPE",
             "Mean Absolute Percentage Error measures the average percentage difference between actual and predicted prices.")
    add_card(slide, 8.75, 1.25, 3.55, 2.15, "R² Score",
             "Measures how much variation is explained by the model. A value nearer to 1 is generally better.",
             RGBColor(141, 91, 166))
    add_rect(slide, 4.8, 3.85, 7.5, 2.3, NAVY)
    add_text(slide, 5.15, 4.15, 6.8, 1.55,
             "Honest viva point:\nGenerated fallback data demonstrates the ML pipeline, but verified local market data is required to prove real-world accuracy.",
             18, WHITE, True, PP_ALIGN.CENTER)

    # 9. Features
    slide = base_slide(prs, "Dashboard Features", 9)
    features = [
        "Crop, region, currency, and forecast-period filters",
        "Current benchmark quote or labelled live-signal estimate",
        "Up to 10 years of historical price data",
        "24-month monthly price forecast",
        "Confidence, model type, observations, and holdout MAPE",
        "Main positive and negative price drivers",
        "Data-source and quality information",
        "CSV and JSON export",
    ]
    add_bullets(slide, 0.75, 1.25, 7.0, 5.4, features, 17)
    add_rect(slide, 8.15, 1.35, 4.35, 4.9, NAVY)
    add_text(slide, 8.55, 1.75, 3.55, 0.5, "Supported Scope", 23, GOLD, True,
             PP_ALIGN.CENTER)
    add_text(slide, 8.55, 2.65, 3.55, 2.8,
             "22\nCROPS & VARIANTS\n\n14\nGLOBAL REGIONS\n\n4\nCURRENCIES",
             20, WHITE, True, PP_ALIGN.CENTER)

    # 10. Demo
    slide = base_slide(prs, "Practical Demonstration", 10)
    add_rect(slide, 0.7, 1.2, 5.15, 4.95, NAVY)
    add_text(slide, 1.0, 1.55, 4.55, 0.45, "Run the application", 22, GOLD, True,
             PP_ALIGN.CENTER)
    add_text(slide, 1.05, 2.25, 4.45, 0.75, "python server.py 8100",
             20, WHITE, True, PP_ALIGN.CENTER, font="Consolas")
    add_text(slide, 1.05, 3.05, 4.45, 0.75, "http://localhost:8100/",
             18, LIGHT_GREEN, True, PP_ALIGN.CENTER, font="Consolas")
    add_text(slide, 1.15, 4.2, 4.25, 1.15,
             "Use Wheat or Rice - Basmati\nIndia • INR • 6 or 12 months",
             17, WHITE, True, PP_ALIGN.CENTER)
    add_bullets(slide, 6.35, 1.35, 5.9, 4.9, [
        "Explain the selected input values.",
        "Show today's price and its source label.",
        "Show historical and future charts.",
        "Point out model type, confidence, and MAPE.",
        "Explain the strongest price drivers.",
        "Export the result as CSV or JSON.",
        "Keep screenshots ready if internet access fails.",
    ], 17)

    # 11. Limitations
    slide = base_slide(prs, "Limitations and Future Scope", 11)
    add_card(slide, 0.7, 1.25, 5.8, 4.95, "Current Limitations",
             "• Public futures are benchmark prices, not exact mandi prices.\n\n"
             "• Not every crop has an observed public historical series.\n\n"
             "• External APIs may be delayed or unavailable.\n\n"
             "• Long-range predictions contain more uncertainty.\n\n"
             "• Generated fallback data cannot prove real-market accuracy.",
             RGBColor(184, 72, 72))
    add_card(slide, 6.85, 1.25, 5.8, 4.95, "Future Enhancements",
             "• Add official Agmarknet and APMC datasets.\n\n"
             "• Create district, mandi, crop-grade, and unit-specific models.\n\n"
             "• Add prediction intervals and rolling backtests.\n\n"
             "• Add login, watchlists, alerts, and maps.\n\n"
             "• Deploy with a database and scheduled model refresh.")

    # 12. Viva
    slide = base_slide(prs, "Key Viva Answers", 12)
    questions = [
        ("Project type?", "A regression and time-series forecasting web application."),
        ("Why ridge regression?", "It is fast, interpretable, handles correlated features, and reduces overfitting."),
        ("What is a lag?", "A previous price value used as an input for forecasting."),
        ("What is MAPE?", "Average prediction error expressed as a percentage; lower is better."),
        ("Is the price guaranteed?", "No. It is a planning estimate affected by market uncertainty."),
        ("Main contribution?", "Current signals, history, forecast, explanation, validation, and export in one dashboard."),
    ]
    for i, (question, answer) in enumerate(questions):
        x = 0.65 + (i % 2) * 6.15
        y = 1.18 + (i // 2) * 1.73
        add_card(slide, x, y, 5.7, 1.35, question, answer,
                 GREEN if i % 2 == 0 else RGBColor(39, 105, 156))

    # 13. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_rect(slide, 0, 0, 13.333, 0.18, GREEN, radius=False)
    add_text(slide, 1.0, 0.85, 11.3, 0.65, "Conclusion", 34, GOLD, True,
             PP_ALIGN.CENTER)
    add_text(slide, 1.4, 1.85, 10.5, 2.45,
             "CropSmart brings together public market signals,\n"
             "historical benchmark data, transparent forecasting,\n"
             "validation information, and reusable exports.",
             25, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 1.55, 4.65, 10.2, 0.7,
             "It is an academic decision-support system, not a guaranteed market oracle.",
             18, LIGHT_GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, 1.0, 5.8, 11.3, 0.75, "THANK YOU", 34, GOLD, True,
             PP_ALIGN.CENTER)
    add_text(slide, 1.0, 6.55, 11.3, 0.35, "Questions?", 18, WHITE, False,
             PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    return len(prs.slides)


if __name__ == "__main__":
    count = build_deck()
    print(f"Created {OUTPUT} with {count} slides")
